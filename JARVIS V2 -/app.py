from flask import Flask, request, jsonify, render_template, g
import requests
import sqlite3
import threading
import os
import sys
import re
import time
import json
import uuid
import shutil
import smtplib
import logging
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as docxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
import google.generativeai as genai

# 🔑 API CONFIG
GEMINI_KEY = ""
GROQ_KEY = ""

# ✉️ EMAIL CONFIG
SENDER_EMAIL = "omegaisgay6969@gmail.com"  # Your Gmail address
SENDER_PASSWORD = ""  # Your Gmail App Password (NOT your regular password)

genai.configure(api_key=GEMINI_KEY)

# ---------- logging ----------
logging.basicConfig(
    level=logging.INFO,
    format="[%(asctime)s] [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger("jarvis")

# -------- SESSION & PENDING OPERATIONS --------
pending_operations = {}  # Store pending confirmations: {session_id: operation_data}
pending_lock = threading.Lock()

# -------- FAST INTERNET CACHE (Restored) --------
_last_check = 0
_cached_status = False

def internet_available():
    global _last_check, _cached_status
    if time.time() - _last_check < 5:
        return _cached_status
    try:
        requests.get("https://www.google.com", timeout=2)
        _cached_status = True
    except Exception as e:
        logger.warning(f"Internet check failed: {e}")
        _cached_status = False
    _last_check = time.time()
    return _cached_status

# -------- DATABASE & RESOURCE UTILS (Restored) --------
DATABASE_PATH = "memory.db"
def resource_path(p):
    try:
        base = sys._MEIPASS
    except Exception:
        base = os.path.abspath(".")
    return os.path.join(base, p)

def get_db():
    if "db" not in g:
        conn = sqlite3.connect(DATABASE_PATH, check_same_thread=False)
        conn.execute("CREATE TABLE IF NOT EXISTS chats (user TEXT, bot TEXT)")
        conn.row_factory = sqlite3.Row
        g.db = conn
    return g.db

def close_db(e=None):
    db = g.pop("db", None)
    if db is not None:
        db.close()

def get_chat_history(limit=50):
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT user, bot FROM chats ORDER BY ROWID DESC LIMIT ?", (limit,))
    rows = cursor.fetchall()
    history = []
    for user_msg, bot_msg in reversed(rows):
        history.append({"role": "user", "content": user_msg})
        history.append({"role": "assistant", "content": bot_msg})
    return history

def clear_chat_history():
    """Clear all chat history from database"""
    try:
        conn = get_db()
        conn.execute("DELETE FROM chats")
        conn.commit()
        return True
    except Exception as e:
        logger.error(f"clear_chat_history failed: {e}")
        return False

# -------- JARVIS IDENTITY (Original) --------
profile = """
You are Jarvis, AI assistant for Sahil. Always address him as Sir.
Sahil is a 4th sem Computer Science student.

MEMORY & CONTEXT:
- You have access to the COMPLETE conversation history with Sir
- Remember ALL previous requests Sir made (documents created, presentations made, files created, etc)
- Reference past interactions when relevant: "As you asked me to create earlier..." or "Like the DBMS doc you requested..."
- Be aware of patterns in Sir's requests and anticipate his needs
- Track what documents/files/presentations you've made for Sir
- Use this context to improve recommendations and understanding

INSTRUCTIONS:
- Do NOT return JSON unless specifically requested for file creation
- Keep chat replies short and professional
- Be proactive: mention related previous work if applicable
- Show you remember Sir's academic journey and work patterns
"""

# -------- UTILS ----------
CREATED_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "created_files")
os.makedirs(CREATED_DIR, exist_ok=True)

def safe_file_name(name: str) -> str:
    name = (name or "").strip()
    if not name:
        return None
    name = re.sub(r"[<>:\"/\\|?*]", "_", name)
    name = re.sub(r"[^\w\-. ]", "_", name)
    name = name.strip()
    return name[:255] if len(name) > 255 else name

def safe_path_under_created(filename: str) -> str:
    if not filename:
        raise ValueError("Filename is required")
    fn = safe_file_name(filename)
    abs_path = os.path.abspath(os.path.join(CREATED_DIR, fn))
    if not abs_path.startswith(os.path.abspath(CREATED_DIR)):
        raise ValueError("Disallowed path")
    return abs_path

def valid_email(email: str) -> bool:
    return bool(re.match(r"^[^@]+@[^@]+\.[^@]+$", (email or "").strip()))

# -------- UNIFIED BRAIN LOGIC (Memory Enabled) --------
def fallback_mistral(system_prompt, history, current_user_msg):
    try:
        messages = [{"role": "system", "content": system_prompt}]
        messages.extend(history)
        messages.append({"role": "user", "content": current_user_msg})
        url = "http://localhost:11434/api/chat"
        payload = {"model": "mistral", "messages": messages, "stream": False}
        response = requests.post(url, json=payload, timeout=20)
        bot_reply = response.json().get('message', {}).get('content', "")
        bot_reply = bot_reply.replace('```json', '').replace('```', '').strip()
        return bot_reply, "Local Mistral 🧠"
    except Exception as e:
        logger.error(f"fallback_mistral failed: {e}")
        return "I am disconnected from all brain modules, Sir.", "Offline"

def ask_ai(current_prompt, is_json=False):
    history = get_chat_history()
    sys_instr = profile
    if is_json:
        sys_instr += "\nSTRICT: Respond ONLY with valid JSON. Create a 'clean_title'. Expand 'body' for word counts."

    if internet_available():
        try:
            model = genai.GenerativeModel("gemini-2.5-flash")
            chat_session = model.start_chat(history=[
                {"role": "user" if m["role"] == "user" else "model", "parts": [m["content"]]}
                for m in history
            ])
            response = chat_session.send_message(f"{sys_instr}\n\nUser: {current_prompt}")
            text = response.text.replace('```json', '').replace('```', '').strip()
            return text, "Gemini ⚡"
        except Exception as e:
            logger.warning(f"Gemini failed: {e}")

        try:
            url = "https://api.groq.com/openai/v1/chat/completions"
            headers = {"Authorization": f"Bearer {GROQ_KEY}", "Content-Type": "application/json"}
            messages = [{"role": "system", "content": sys_instr}]
            messages.extend(history)
            messages.append({"role": "user", "content": current_prompt})
            payload = {"model": "llama-3.3-70b-versatile", "messages": messages, "temperature": 0.7}
            res = requests.post(url, json=payload, headers=headers, timeout=10)
            text = res.json().get('choices', [])[0].get('message', {}).get('content', "")
            text = text.replace('```json', '').replace('```', '').strip()
            return text, "Groq 🚀"
        except Exception as e:
            logger.warning(f"Groq failed: {e}")

    return fallback_mistral(sys_instr, history, current_prompt)

def get_recently_created_files(limit=5):
    """Get list of recently created documents and presentations"""
    try:
        if not os.path.exists(CREATED_DIR):
            return []
        files = sorted(
            [f for f in os.listdir(CREATED_DIR) if f.endswith(('.docx', '.pptx'))],
            key=lambda x: os.path.getmtime(os.path.join(CREATED_DIR, x)),
            reverse=True
        )
        return files[:limit]
    except Exception as e:
        logger.error(f"get_recently_created_files failed: {e}")
        return []

# -------- FILE ENGINES (Professional Positioning) --------
def create_perfect_ppt(data):
    prs = Presentation()
    clean_title = data.get('clean_title', 'Presentation')
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = clean_title.upper()
    slide.placeholders[1].text = "Technical Briefing // Prepared for Sahil Sir"

    for item in data.get('slides', []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        t_p = t_box.text_frame.paragraphs[0]
        t_p.text = item.get('title', '')
        t_p.font.name, t_p.font.size, t_p.font.bold = 'Segoe UI', Pt(32), True
        b_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5), Inches(5))
        tf = b_box.text_frame
        tf.word_wrap = True
        for point in item.get('content', []):
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.name, p.font.size = 'Segoe UI', Pt(18)
            p.space_after = Pt(10)
        try:
            keyword = item.get('image_keyword', 'tech')
            img_data = requests.get(f"https://loremflickr.com/800/600/{keyword}", timeout=5).content
            temp_file = os.path.join(CREATED_DIR, f"temp_{uuid.uuid4().hex}.jpg")
            with open(temp_file, "wb") as f:
                f.write(img_data)
            slide.shapes.add_picture(temp_file, Inches(5.8), Inches(1.3), width=Inches(3.7))
            os.remove(temp_file)
        except Exception as e:
            logger.warning(f"create_perfect_ppt image fallback: {e}")

    file_name = safe_file_name(clean_title) or "Presentation"
    path = os.path.join(CREATED_DIR, f"{file_name}.pptx")
    prs.save(path)
    return path

def create_perfect_doc(data):
    doc = Document()
    clean_title = data.get('clean_title', 'Document')
    doc.add_heading(clean_title.upper(), 0).alignment = WD_ALIGN_PARAGRAPH.CENTER
    for sec in data.get('sections', []):
        h = doc.add_heading(sec.get('heading', ''), level=1)
        h.runs[0].font.name, h.runs[0].font.size = 'Segoe UI', docxPt(16)
        p = doc.add_paragraph(sec.get('body', ''))
        p.paragraph_format.line_spacing = 1.15
        p.runs[0].font.name, p.runs[0].font.size = 'Segoe UI', docxPt(11)

    file_name = safe_file_name(clean_title) or "Document"
    path = os.path.join(CREATED_DIR, f"{file_name}.docx")
    doc.save(path)
    return path

# -------- OPERATION DETECTION & CONFIRMATION SYSTEM --------
def detect_operation_intent(user_input):
    lower = user_input.lower()
    if re.search(r'\b(create|make|write|edit|add|delete|remove|open|read)\b.*\b(file|txt|text|note|notepad|document|content)\b', lower):
        return "file_operation"
    if re.search(r'\b(create|make|new|mkdir|build)\b.*\b(folder|directory|dir)\b', lower):
        return "folder_operation"
    if re.search(r'\b(send|mail|email)\b', lower):
        return "email_operation"
    if re.search(r'\b(make|create|generate|prepare|build|new)\b.*\b(ppt|presentation|slideshow)\b', lower) or lower.startswith("ppt on"):
        return "ppt_operation"
    if re.search(r'\b(make|create|generate|prepare|write|build|new)\b.*\b(doc|essay|report|document|word|docx)\b', lower) or lower.startswith(("doc on", "essay on")):
        return "doc_operation"
    if re.search(r'\b(clear|delete|wipe|erase|reset|remove)\b.*\b(memory|chat|history|record)\b', lower):
        return "clear_memory_operation"
    if re.search(r'\b(forget|forget everything|clear everything)\b', lower):
        return "clear_memory_operation"
    return None

def parse_operation_request(raw_input, operation_type):
    if operation_type == "file_operation":
        prompt = f"""You are a file operation interpreter. Analyze this user request carefully and extract the INTENT:
User Request: '{raw_input}'

Determine:
1. What file action does user want? (create new file, edit existing, read contents, delete file, append content)
2. What should the file be named? (use .txt extension if not specified)
3. What content should go in the file? (extract ALL content the user wants in the file, not just first words)
4. Is user referencing an existing file or creating new?

Return ONLY valid JSON with all fields filled based on your understanding:
{{"operation": "file_operation", "action": "create|edit|delete|read|append", "filename": "name.ext", "content": "full content here"}}"""
    elif operation_type == "folder_operation":
        prompt = f"""You are a folder operation interpreter. Analyze this user request:
User Request: '{raw_input}'

Determine:
1. Should this create or delete a folder?
2. What should the folder be named? (make it descriptive)
3. What path should it be created in? (if not specified, use current directory)

Return ONLY valid JSON:
{{"operation": "folder_operation", "action": "create|delete", "folder_path": "appropriate/path/or/name"}}"""
    elif operation_type == "email_operation":
        prompt = f"""You are an email operation interpreter. Analyze this user request carefully:
User Request: '{raw_input}'

Determine:
1. Who is the recipient? (extract email address or determine from context)
2. What should the subject line be? (make it concise and meaningful)
3. What is the email body? (extract the full message the user wants to send)
4. Is this a professional, casual, or specific type of email?

Return ONLY valid JSON with complete details:
{{"operation": "email_operation", "recipient": "user@email.com", "subject": "appropriate subject line", "body": "complete email message body here"}}"""
    elif operation_type == "ppt_operation":
        prompt = f"""You are a presentation generator. Analyze this user request:
User Request: '{raw_input}'

Understand:
1. What is the main topic? (extract the core subject)
2. How many slides are needed? (default 5 if not specified)
3. What level of detail? (academic, professional, beginner, expert)
4. Any specific requirements? (themes, focus areas, examples needed?)

Return ONLY valid JSON:
{{"operation": "ppt_operation", "topic": "clear topic title", "num_slides": 5, "level": "academic/professional/beginner", "details": "specific requirements or focus areas"}}"""
    elif operation_type == "doc_operation":
        prompt = f"""You are a document generator for academic/professional content. Analyze this user request:
User Request: '{raw_input}'

Understand:
1. What is the main topic/subject?
2. How comprehensive should it be? (extract word count if mentioned, else estimate)
3. What type of document? (notes, essay, report, tutorial, guide)
4. What level? (beginner, intermediate, advanced, academic)
5. Should it include examples, questions, worked solutions?

Return ONLY valid JSON:
{{"operation": "doc_operation", "topic": "clear topic", "type": "notes|essay|report|guide", "level": "academic|professional|beginner", "include_examples": true, "include_questions": true, "estimated_words": 5000}}"""
    else:
        return None

    res, _ = ask_ai(prompt, is_json=True)
    try:
        parsed = json.loads(res)
        parsed["operation"] = operation_type
        return parsed
    except Exception as e:
        logger.warning(f"parse_operation_request failed: {e} ; res={res}")
        return None

def create_confirmation_message(details):
    op_type = details.get("operation")
    if op_type == "file_operation":
        action = details.get("action", "").upper()
        filename = details.get("filename", "file")
        content_preview = details.get("content", "")[:50]
        if len(details.get("content", "")) > 50:
            content_preview += "..."
        msg = f"📄 File Operation:\n- Action: {action}\n- Filename: {filename}"
        if content_preview and action in ["CREATE", "EDIT"]:
            msg += f"\n- Content Preview: {content_preview}\n\nProceed? (yes/no)"
        else:
            msg += "\n\nProceed? (yes/no)"
        return msg
    if op_type == "folder_operation":
        action = details.get("action", "").upper()
        folder = details.get("folder_path", "folder")
        return f"📁 Folder Operation:\n- Action: {action}\n- Path: {folder}\n\nProceed? (yes/no)"
    if op_type == "email_operation":
        recipient = details.get("recipient", "recipient")
        subject = details.get("subject", "subject")
        return f"📧 Email Operation:\n- To: {recipient}\n- Subject: {subject}\n\nProceed? (yes/no)"
    if op_type == "ppt_operation":
        topic = details.get("topic", "presentation")
        return f"🎫 PowerPoint Creation:\n- Topic: {topic}\n- Will create and open automatically\n\nProceed? (yes/no)"
    if op_type == "doc_operation":
        topic = details.get("topic", "document")
        return f"📝 Document Creation:\n- Topic: {topic}\n- Will create and open automatically\n\nProceed? (yes/no)"
    if op_type == "clear_memory_operation":
        return f"🧠 Memory Clear Operation:\n- Action: Delete all chat history and memory\n- This cannot be undone\n\nProceed? (yes/no)"
    return "Proceed with this operation? (yes/no)"

def execute_file_operation(details):
    action = details.get("action", "").lower()
    filename = details.get("filename", "file.txt")
    content = details.get("content", "")
    filename_safe = safe_file_name(filename)
    if not filename_safe:
        return "❌ Invalid filename, Sir."
    filepath = safe_path_under_created(filename_safe)

    try:
        if action == "create":
            os.makedirs(os.path.dirname(filepath) or CREATED_DIR, exist_ok=True)
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
            return f"✅ File '{filename_safe}' created successfully, Sir. Location: {filepath}"
        elif action == "edit":
            with open(filepath, "a", encoding="utf-8") as f:
                f.write("\n" + content)
            return f"✅ File '{filename_safe}' updated successfully, Sir."
        elif action == "read":
            with open(filepath, "r", encoding="utf-8") as f:
                file_content = f.read()
            return f"📖 File contents:\n\n{file_content}"
        elif action == "delete":
            os.remove(filepath)
            return f"✅ File '{filename_safe}' deleted successfully, Sir."
        else:
            return f"❌ Unknown file action: {action}"
    except Exception as e:
        logger.error(f"execute_file_operation failed: {e}")
        return f"❌ Error executing file operation: {str(e)}"

def execute_folder_operation(details):
    action = details.get("action", "").lower()
    folder_path = details.get("folder_path", "new_folder")
    folder_path = os.path.normpath(folder_path)
    if not os.path.isabs(folder_path):
        folder_path = os.path.join(CREATED_DIR, folder_path)

    try:
        if action == "create":
            os.makedirs(folder_path, exist_ok=True)
            return f"✅ Folder '{folder_path}' created successfully, Sir."
        elif action == "delete":
            shutil.rmtree(folder_path)
            return f"✅ Folder '{folder_path}' deleted successfully, Sir."
        else:
            return f"❌ Unknown folder action: {action}"
    except Exception as e:
        logger.error(f"execute_folder_operation failed: {e}")
        return f"❌ Error executing folder operation: {str(e)}"

def format_email_html(subject, content):
    return f"""
    <html>
    <body style="font-family:Segoe UI,Arial,sans-serif; line-height:1.6; color:#333;">
        
        <div style="max-width:600px; margin:auto; padding:20px; border:1px solid #ddd; border-radius:10px;">
            
            <h2 style="color:#2c3e50;">{subject}</h2>
            
            <div style="margin-top:15px;">
                {content}
            </div>

            <br>

        </div>

    </body>
    </html>
    """

def execute_email_operation(details):
    recipient = details.get("recipient", "").strip()
    subject = details.get("subject", "").strip()
    raw_body = details.get("body", "").strip()

    if not valid_email(recipient):
        return f"❌ Invalid email address: {recipient}. Please provide a valid email."
    if not SENDER_EMAIL or not SENDER_PASSWORD:
        return "❌ Email credentials not configured, Sir."

    try:
        subject_ai, _ = ask_ai(f"""
        Rewrite this as a short professional email subject.

        RULES:
        - Return ONLY one subject line
        - No explanations
        - Max 10 words
        - content should be concise and reflect the email body

        Input: {subject}
        """)
        subject = subject_ai.split("\n")[0].strip() if subject_ai else subject

        ai_body, _ = ask_ai(f"""
Write a professional email in HTML BODY format.

STRICT RULES:
- Output ONLY HTML (no explanations)
- NO ```html or markdown
- NO "Would you like..." or suggestions
- NO extra commentary
- Start with <p>Dear Sir,</p>
- Use <p> for paragraphs
- Keep it around 50 words

CONTENT:
{raw_body}
""")
        ai_body_filtered = ai_body.replace("```html", "").replace("```", "").strip()
        html_body = format_email_html(subject, ai_body_filtered)

        msg = MIMEMultipart("alternative")
        msg['From'] = SENDER_EMAIL
        msg['To'] = recipient
        msg['Subject'] = subject
        msg.attach(MIMEText(html_body, 'html'))

        server = smtplib.SMTP('smtp.gmail.com', 587, timeout=20)
        server.starttls()
        server.login(SENDER_EMAIL, SENDER_PASSWORD)
        server.send_message(msg)
        server.quit()

        return f"✅ Smart formatted email sent to {recipient}, Sir."
    except Exception as e:
        logger.error(f"execute_email_operation failed: {e}")
        return f"❌ Error sending email: {str(e)}"

def process_pending(op_data):
    op_type = op_data.get("operation")
    if op_type == "file_operation":
        return execute_file_operation(op_data), "File System 📁"
    elif op_type == "folder_operation":
        return execute_folder_operation(op_data), "File System 📁"
    elif op_type == "email_operation":
        return execute_email_operation(op_data), "Email Client 📧"
    elif op_type == "ppt_operation":
        topic = op_data.get('topic', op_data.get('original_request'))
        num_slides = op_data.get('num_slides', 5)
        level = op_data.get('level', 'professional')
        details = op_data.get('details', '')
        prompt = f"""Generate a professional {num_slides}-slide presentation on: '{topic}'

PRESENTATION SPECIFICATIONS:
- Number of slides: {num_slides} (after title slide)
- Level: {level} (for {level} audience)
- Special requirements: {details if details else 'Standard comprehensive presentation'}

SLIDE STRUCTURE:
1. Each slide MUST have a clear title
2. Content should be organized as bullet points with detailed explanations
3. Include practical examples and key takeaways
4. Make slides visually oriented (suggest relevant image keywords)
5. Progress logically from introduction to conclusion

CONTENT REQUIREMENTS:
- Provide substantial, meaningful content for each slide
- Include examples where applicable
- Ensure every slide adds educational value
- Make it engaging and informative

Return ONLY valid JSON - NO markdown:
{{'clean_title': 'Presentation Title', 'slides': [{{'title': 'Slide Title', 'content': ['Point 1 with explanation', 'Point 2 with explanation', 'Key takeaway or example'], 'image_keyword': 'relevant-keyword-for-image'}}]}}"""
        res, model = ask_ai(prompt, is_json=True)
        try:
            data = json.loads(res)
            file_path = create_perfect_ppt(data)
            os.startfile(file_path)
            return f"✅ The presentation on {data.get('clean_title', 'Untitled')} is ready, Sir.", model
        except Exception as e:
            logger.error(f"process_pending ppt failed: {e}")
            return f"❌ I was unable to format the PPT data correctly, Sir. Error: {str(e)}", "Jarvis 🤖"
    elif op_type == "doc_operation":
        original = op_data.get('original_request')
        doc_type = op_data.get('type', 'notes')
        doc_level = op_data.get('level', 'academic')
        include_examples = op_data.get('include_examples', True)
        include_questions = op_data.get('include_questions', True)
        word_target = op_data.get('estimated_words', 5000)
        prompt = f"""Generate a comprehensive {doc_type} document for: '{original}'

DOCUMENT SPECIFICATIONS:
- Type: {doc_type} (academic notes with depth and clarity)
- Level: {doc_level} (suitable for {doc_level} audience)
- Target length: approximately {word_target} words
- Include detailed explanations and key concepts
- Include worked examples: {include_examples}
- Include practice questions/problems: {include_questions}

CONTENT REQUIREMENTS:
- Start with introduction and overview
- Create multiple sections with clear headings
- Each section should be substantial and informative
- Include real-world examples where applicable
- If questions included: provide solutions too
- End with summary of key points

Return ONLY valid JSON - NO markdown:
{{'clean_title': 'Document Title', 'sections': [{{'heading': 'Section Name', 'body': 'Detailed content here with full explanations, examples, and if applicable, worked solutions...'}}]}}"""
        res, model = ask_ai(prompt, is_json=True)
        try:
            data = json.loads(res)
            file_path = create_perfect_doc(data)
            os.startfile(file_path)
            return f"✅ Document '{data.get('clean_title', 'Untitled')}' is ready, Sir.", model
        except Exception as e:
            logger.error(f"process_pending doc failed: {e}")
            return f"❌ Document compilation error, Sir. Error: {str(e)}", "Jarvis 🤖"
    elif op_type == "clear_memory_operation":
        if clear_chat_history():
            return "✅ All chat history and memory cleared successfully, Sir. Starting fresh!", "Memory System 🧠"
        return "❌ Error clearing memory, Sir.", "Memory System 🧠"
    else:
        return "❌ Unknown operation type, Sir.", "Jarvis 🤖"

# -------- FLASK ROUTES --------
app = Flask(__name__, template_folder=resource_path("templates"))
app.config['SECRET_KEY'] = 'jarvis-secret-' + str(uuid.uuid4())
app.teardown_appcontext(close_db)

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/status")
def status():
    conn = get_db()
    count = conn.execute("SELECT COUNT(*) FROM chats").fetchone()[0]
    return jsonify({"internet": internet_available(), "model": "Jarvis Memory Core", "memories": count})

@app.route("/history")
def history():
    conn = get_db()
    rows = conn.execute("SELECT user, bot FROM chats ORDER BY ROWID DESC LIMIT 50").fetchall()
    return jsonify([{"user": r["user"], "bot": r["bot"]} for r in reversed(rows)])

@app.route("/files")
def list_files():
    try:
        files = sorted(
            [f for f in os.listdir(CREATED_DIR) if f.endswith(('.docx', '.pptx'))],
            key=lambda x: os.path.getmtime(os.path.join(CREATED_DIR, x)),
            reverse=True
        )
        return jsonify(files[:8])
    except Exception as e:
        logger.error(f"list_files failed: {e}")
        return jsonify([])

@app.route("/open_local/<filename>")
def open_local(filename):
    try:
        safe_name = safe_file_name(filename)
        if not safe_name: raise ValueError("Invalid filename")
        filepath = os.path.join(CREATED_DIR, safe_name)
        if os.path.exists(filepath):
            os.startfile(filepath)
            return jsonify({"status": "opened"})
        return jsonify({"status": "not found"}), 404
    except Exception as e:
        logger.error(f"open_local failed: {e}")
        return jsonify({"status": "error", "error": str(e)}), 500

@app.route("/welcome")
def welcome():
    res, _ = ask_ai("Give a very short greeting to Sahil Sir.")
    return jsonify({"reply": res})

@app.route("/chat", methods=["POST"])
def chat():
    raw = request.json.get("message", "")
    user_input = raw.lower().strip()
    session_id = request.json.get("session_id", str(uuid.uuid4()))
    model = "Jarvis 🤖"

    if user_input in ["yes", "no"]:
        with pending_lock:
            op_data = pending_operations.pop(session_id, None)
            if not op_data and pending_operations:
                first_session = next(iter(pending_operations))
                op_data = pending_operations.pop(first_session, None)

        if op_data:
            if user_input == "yes":
                reply, model = process_pending(op_data)
            else:
                reply = "❌ Operation cancelled, Sir."
        else:
            reply, model = ask_ai(raw)
    else:
        op_intent = detect_operation_intent(user_input)
        if op_intent:
            op_details = parse_operation_request(raw, op_intent)
            if op_details:
                op_details["original_request"] = raw
                with pending_lock:
                    pending_operations[session_id] = op_details
                reply = create_confirmation_message(op_details)
                model = "Operation Parser 🔍"
            else:
                reply = "❌ I couldn't parse that operation, Sir. Could you please rephrase?"
                model = "Operation Parser 🔍"
        else:
            reply, model = ask_ai(raw)

    try:
        conn = get_db()
        conn.execute("INSERT INTO chats VALUES (?, ?)", (raw, reply))
        conn.commit()
    except Exception as e:
        logger.error(f"save chat failed: {e}")

    return jsonify({"reply": reply, "session_id": session_id, "model": model})

if __name__ == "__main__":
    chrome_cmd = 'start chrome --app=http://127.0.0.1:5000 --window-size=1000,700 --window-position=100,100'
    threading.Timer(2.0, lambda: os.system(chrome_cmd)).start()
    app.run(debug=False, port=5000)